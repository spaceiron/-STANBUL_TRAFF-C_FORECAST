"""
Generates English PPT for Software Requirements & Analysis class.
Output: Istanbul_Traffic_Forecast_Requirements_Analysis_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

OUTPUT = "Istanbul_Traffic_Forecast_Requirements_Analysis_Presentation.pptx"


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle
    return slide


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(28)
    tx.text_frame.paragraphs[0].font.bold = True

    def box(x, y, w, h, heading, items):
        t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = t.text_frame
        tf.text = heading
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(20)
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(14)
        return t

    box(0.5, 1.1, 4.5, 5.5, left_title, left_bullets)
    box(5.2, 1.1, 4.3, 5.5, right_title, right_bullets)
    return slide


def add_diagram_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.55))
    t.text_frame.text = title
    t.text_frame.paragraphs[0].font.size = Pt(26)
    t.text_frame.paragraphs[0].font.bold = True
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    tf.text = lines[0]
    tf.paragraphs[0].font.name = "Courier New"
    tf.paragraphs[0].font.size = Pt(11)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Courier New"
        p.font.size = Pt(11)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Software Requirements & Analysis\nProject Presentation",
        "Istanbul Traffic Forecast (TrafikPuls)\nUzay Demir | Student No: 210218027\n"
        "Department of Computer Engineering, Istanbul Okan University",
    )

    add_bullet_slide(
        prs,
        "Motivation",
        [
            "Public transport occupancy information is fragmented and hard to trust in real time.",
            "Combining ML prediction with passenger crowdsourcing creates a feedback loop.",
            "Goal: better route decisions, early density alerts, and transparent confidence signals.",
        ],
    )

    add_bullet_slide(
        prs,
        "Project Vision",
        [
            "A mobile service that shows density forecasts for favorite lines, explains reliability,",
            "suggests less crowded alternatives when possible, and learns from passenger reports.",
            "Near-term: working prototype (map, prediction, feedback, profile, admin).",
            "Mid-term: validated open-data integration and production-grade push triggers.",
        ],
    )

    add_two_column_slide(
        prs,
        "Functional vs Non-Functional (one example each)",
        "Functional requirement",
        [
            "For a selected route, the system shall return:",
            "• density score + confidence",
            "• stop list with ETA",
            "• 30/60 min forecast via Flask API",
            "User shall submit occupancy feedback (empty / standing / full)",
            "and the backend shall blend recent reports into predictions.",
        ],
        "Non-functional requirement",
        [
            "Availability / resilience:",
            "If the live vehicle feed is unavailable, the API shall fall back",
            "to mock vehicles so predictions still respond.",
            "Typical API calls shall complete within ~10 seconds under normal load.",
        ],
    )

    add_bullet_slide(
        prs,
        "Requirement Elicitation",
        [
            "Review of SRS/SDD scenarios (map search, prediction sheet, reporting).",
            "Benchmarking similar mobility apps and open-data portals.",
            "User journey drafting (commute, favorite line, threshold alert).",
            "Technical constraints workshop (mobile + HTTPS + Firebase rules).",
        ],
    )

    add_bullet_slide(
        prs,
        "Requirement Engineering — Sub-topics",
        [
            "1) Elicitation — discovering needs from stakeholders and documents.",
            "2) Specification — writing SRS/SDD: functional, data, interface, quality attributes.",
            "3) Validation — “Are we building the right product?” (stakeholder fit, acceptance tests).",
            "Related: Verification — “Did we build it right?” (tests, health checks, demos).",
        ],
    )

    add_bullet_slide(
        prs,
        "Stakeholders",
        [
            "Passengers / app users — predictions, ETA, language, alerts.",
            "Data providers (e.g., municipal open data) — feed quality and access policy.",
            "University / course staff — grading, documentation, demonstration.",
            "Developers / maintainers — architecture, deployment, backlog.",
            "Operators (optional) — monitoring suspicious feedback, model operations.",
        ],
    )

    add_two_column_slide(
        prs,
        "Validation vs Verification",
        "Validation",
        [
            "Maps user stories to implemented features.",
            "Acceptance-style checks: can a user select M4, see forecast,",
            "submit feedback, and see incidents/alternatives?",
            "Evidence: acceptance test report + live demo.",
        ],
        "Verification",
        [
            "Checks implementation against specification.",
            "Examples: GET /health returns 200 JSON;",
            "prediction JSON contains required fields;",
            "Flutter uses configurable API_BASE_URL for cloud.",
        ],
    )

    add_bullet_slide(
        prs,
        "Prototype",
        [
            "End-to-end working prototype on real stack:",
            "Flutter UI, Firebase Auth/Firestore, Flask ML API (Render HTTPS).",
            "Purpose: reduce uncertainty early (API latency, map UX, feedback loop).",
        ],
    )

    add_bullet_slide(
        prs,
        "Why a low-fidelity / iterative prototype approach?",
        [
            "Low cost of change: prioritize flows and contracts before pixel-perfect UI.",
            "Technical risk first: map + network + model integration is harder than static mockups.",
            "Course timeline: functional skeleton validates requirements faster than high-fi Figma only.",
            "Note: UI is increasingly polished, but the method remains iterative and evidence-driven.",
        ],
    )

    add_bullet_slide(
        prs,
        "Requirement Management",
        [
            "Version control (Git/GitHub) for traceability of changes.",
            "Link SRS items to features (e.g., threshold alerts ↔ settings service).",
            "Backlog for deferred items (production FCM server push, full RBAC admin).",
            "Deployment artifacts tracked (Procfile, requirements.txt, render.yaml).",
        ],
    )

    add_diagram_slide(
        prs,
        "Context Diagram (textual)",
        [
            "   [Passenger] ----uses----> [Flutter App]",
            "        |                          |",
            "        |                          +----> [Firebase: Auth, Firestore, FCM]",
            "        |                          |",
            "        |                          +----> [Google Maps]",
            "        |                          |",
            "        +--------------------------+----> [Flask API on HTTPS (Render)]",
            "                                     |",
            "                                     +--> [(Optional) Open Data API]",
        ],
    )

    add_diagram_slide(
        prs,
        "Use Cases (summary)",
        [
            "Passenger: Register/Login | Search route | View prediction | View stops+ETA",
            "Passenger: Submit feedback | Manage favorites | Set alert threshold | TR/EN",
            "Passenger: View incidents / alternative suggestions (when shown)",
            "Admin: Monitor feedback stream | Review suspicious patterns",
            "Admin: View model metrics | Trigger prototype retrain / push simulation",
        ],
    )

    add_bullet_slide(
        prs,
        "“Hardware” & cognition-style view (software project)",
        [
            "Sensing: device location (optional), passenger observation (feedback),",
            "vehicle positions from live feed or mock fallback.",
            "Reasoning: Random Forest density model + recency-weighted user reports + confidence.",
            "Actuation: local/FCM notifications, map markers/polylines, dashboard actions.",
        ],
    )

    add_diagram_slide(
        prs,
        "Data Flow (high level)",
        [
            "  Flutter UI  --(HTTPS)-->  Flask /get_prediction, /get_stops, /get_incidents",
            "       |                              |",
            "       +---- Firestore (feedbacks) <--+---- in-memory user reports (prototype)",
            "       |",
            "       +---- Firebase Auth (session)",
            "  Optional: FCM token stored for future server-triggered alerts",
        ],
    )

    add_bullet_slide(
        prs,
        "Functional vs Non-functional — quick recap",
        [
            "Functional: route-based prediction, ETA, feedback, incidents, alternatives.",
            "Non-functional: HTTPS deploy, resilience fallback, responsiveness, i18n.",
        ],
    )

    add_title_slide(
        prs,
        "Thank you",
        "Questions?\nRepository + live API referenced in project documentation.",
    )

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
